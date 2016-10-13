"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html

Example program showing how to read in Excel, process with pandas and
output to a PowerPoint file.
"""

from __future__ import print_function
from pptx import Presentation
from pptx.chart.data import ChartData 
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import os
import argparse
import MySQLdb as mdb
import pandas as pd
import numpy as np
from datetime import date
import matplotlib.pyplot as plt
import seaborn as sns
import shutil as s

sys.path.insert(0,'/home/djohnson/analytics/python_libs');
from common_libs import *
from create_mysql import *

cur_datetime = datetime.now()
pd.set_option('display.width',1000)
pd.set_option('display.max_colwidth',200)
pd.set_option('display.max_rows',400)

def df_to_table(slide, df, left, top, width, height, colnames=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py
     """
    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name

    m = df.as_matrix()

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create ppt report')
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file used as the template')
    parser.add_argument('report',
                        type=argparse.FileType('r'),
                        help='Excel file containing the raw report data')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint report file')
    return parser.parse_args()


def create_pivot(df, index_list=["Manager", "Rep", "Product"],
                 value_list=["Price", "Quantity"]):
    """
    Take a DataFrame and create a pivot table
    Return it as a DataFrame pivot table
    """
    table = pd.pivot_table(df, index=index_list,
                           values=value_list,
                           aggfunc=[np.sum, np.mean], fill_value=0)
    return table


def create_chart(df, filename):
    """ Create a simple bar chart saved to the filename based on the dataframe
    passed to the function
    """
    df['total'] = df['Quantity'] * df['Price']
    final_plot = df.groupby('Name')['total'].sum().order().plot(kind='barh')
    fig = final_plot.get_figure()
    fig.set_size_inches(6, 4.5)
    fig.savefig(filename, bbox_inches='tight', dpi=600)


def create_ppt(input, output, report_data, chart):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Quarterly Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
    # Create the summary graph
    graph_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Sales by account"
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(chart)
    subtitle = slide.placeholders[2]
    subtitle.text = "Results consistent with last quarter"
    # Create a slide for each manager
    for manager in report_data.index.get_level_values(0).unique():
        #print(report_data.xs(manager, level=0).reset_index())
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        title = slide.shapes.title
        title.text = "Report for {}".format(manager)
        top = Inches(1.5)
        left = Inches(0.25)
        width = Inches(9.25)
        height = Inches(5.0)
        # Flatten the pivot table by resetting the index
        # Create a table on the slide
        df_to_table(slide, report_data.xs(manager, level=0).reset_index(),
                    left, top, width, height)
    prs.save(output)


if __name__ == "__main__":

	con = None
	con = mdb.connect('localhost','root','','');
	cur = con.cursor()

	###############
	# Read in data
	###############
	account_id = 1091
	DBNAME = 'benchmark_prod'
	TABLENAME_EFF = 'TMP_EFFECTIVENESS_ALL' 
	TABLENAME_REACH_act = 'AER_REACH_SUMMARY_DATEINTERVAL_account' 
	TABLENAME_REACH_video = 'AER_REACH_SUMMARY_VIDEO_DATEINTERVAL_account' 

	query_eff = "SELECT * FROM %s.%s" % (DBNAME,TABLENAME_EFF)
	query_reach_act = "SELECT * FROM %s.%s WHERE cur_year = '1-Year'" % (DBNAME,TABLENAME_REACH_act)
	query_reach_video = "SELECT * FROM %s.%s WHERE cur_year = '1-Year'" % (DBNAME,TABLENAME_REACH_video)

	#eff_df = createDF_from_MYSQL_query(query_eff)
	reach_act_df = createDF_from_MYSQL_query(query_reach_act)
	reach_video_df = createDF_from_MYSQL_query(query_reach_video)

	#args = parse_args()
	#infile = 'simple-template.pptx'	
	infile = 'qbr_template.pptx'	
	outfile = 'myreport.pptx'
 	report_name = './data/sales-funnel.xlsx' 
	chart = "report-image.png"

	df = pd.read_excel(report_name)
	report_data = create_pivot(df)
	create_chart(df, "report-image.png")

	prs = Presentation(infile)
	# Use the output from analyze_ppt to understand which layouts and placeholders
	# to use
	# Create a title slide first
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	#subtitle = slide.placeholders[1]
	title.text = "Quarterly Report"
	#subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())

	# Create the summary graph
	graph_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(graph_slide_layout)
	title = slide.shapes.title
	title.text = "Yearly Results"
	#placeholder = slide.placeholders[1]
	#pic = placeholder.insert_picture(chart)
	#subtitle = slide.placeholders[2]
	#subtitle.text = "Results consistent with last quarter"
	
	# define chart data ---------------------
	chart_data = ChartData()
	chart_data.categories = ['East', 'West', 'Midwest']
	test = {19.2, 21.4, 16.7}
	chart_data.add_series('Series 1', test )

	# add chart to slide --------------------
	x, y, cx, cy = Inches(3), Inches(1), Inches(6), Inches(4.5)
	slide.shapes.add_chart(
		XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
	)

	# Create a slide for each manager
#	for manager in report_data.index.get_level_values(0).unique():
#		#print(report_data.xs(manager, level=0).reset_index())
#		slide = prs.slides.add_slide(prs.slide_layouts[2])
#		title = slide.shapes.title
#		title.text = "Report for {}".format(manager)
#		top = Inches(1.5)
#		left = Inches(0.25)
#		width = Inches(9.25)
#		height = Inches(5.0)
#		# Flatten the pivot table by resetting the index
#		# Create a table on the slide
#		df_to_table(slide, report_data.xs(manager, level=0).reset_index(),
#                    left, top, width, height)

	try:
		os.remove('/media/sf_transfer/' + outfile)
	except:
		print('No file\n')
	prs.save('/media/sf_transfer/' + outfile)
	#s.copyfile(outfile,'/media/sf_transfer/' + outfile)

	#df = pd.read_excel(args.report.name)
	#report_data = create_pivot(df)
	#create_chart(df, "report-image.png")
	#create_ppt(args.infile.name, args.outfile.name, report_data, "report-image.png")
