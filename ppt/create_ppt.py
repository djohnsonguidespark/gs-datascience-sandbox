#"""
#See http://pbpython.com/creating-powerpoint.html for details on this script
#Requires https://python-pptx.readthedocs.org/en/latest/index.html
#
#Example program showing how to read in Excel, process with pandas and
#output to a PowerPoint file.
#"""

from __future__ import print_function
from pptx import Presentation
from pptx.chart.data import ChartData 
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt 
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_LEGEND_POSITION

import os
import sys 
import argparse
import MySQLdb as mdb
import pandas as pd
import numpy as np
from datetime import date
import matplotlib.pyplot as plt
import seaborn as sns
import shutil as s

sys.path.insert(0,'/home/analytics/analytics_sandbox/python_libs');
from common_libs import *
from create_mysql import *

cur_datetime = datetime.now()
pd.set_option('display.width',1000)
pd.set_option('display.max_colwidth',200)
pd.set_option('display.max_rows',400)

DBNAME = 'benchmark_prod'
SANDBOX = 'sandbox_prod'
TABLENAME_EFF = 'TMP_EFFECTIVENESS_ALL' 
TABLENAME_REACH_act = 'AER_REACH_SUMMARY_DATEINTERVAL_account' 
TABLENAME_REACH_act_2014 = 'AER_REACH_SUMMARY_DATEINTERVAL_account_2014' 
TABLENAME_REACH_video = 'AER_REACH_SUMMARY_VIDEO_DATEINTERVAL_account' 
TABLENAME_REACH_video_2014 = 'AER_REACH_SUMMARY_VIDEO_DATEINTERVAL_account_2014' 
TABLENAME_SURVEY = "TMP_SURVEY_QUESTION_ANSWER_2015OE_Aug"

ADD_DATA_LABELS = False

timeframe_days = 365

def df_to_table(slide, df, left, top, width, height, colnames=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py
     """
    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name

    m = df.as_matrix()

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create ppt report')
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file used as the template')
    parser.add_argument('report',
                        type=argparse.FileType('r'),
                        help='Excel file containing the raw report data')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint report file')
    return parser.parse_args()


def create_pivot(df, index_list=["Manager", "Rep", "Product"],
                 value_list=["Price", "Quantity"]):
    """
    Take a DataFrame and create a pivot table
    Return it as a DataFrame pivot table
    """
    table = pd.pivot_table(df, index=index_list,
                           values=value_list,
                           aggfunc=[np.sum, np.mean], fill_value=0)
    return table


def create_chart(df, filename):
    """ Create a simple bar chart saved to the filename based on the dataframe
    passed to the function
    """
    df['total'] = df['Quantity'] * df['Price']
    final_plot = df.groupby('Name')['total'].sum().order().plot(kind='barh')
    fig = final_plot.get_figure()
    fig.set_size_inches(6, 4.5)
    fig.savefig(filename, bbox_inches='tight', dpi=600)


def create_ppt(input, output, report_data, chart):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Quarterly Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
    # Create the summary graph
    graph_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Sales by account"
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(chart)
    subtitle = slide.placeholders[2]
    subtitle.text = "Results consistent with last quarter"
    # Create a slide for each manager
    for manager in report_data.index.get_level_values(0).unique():
        #print(report_data.xs(manager, level=0).reset_index())
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        title = slide.shapes.title
        title.text = "Report for {}".format(manager)
        top = Inches(1.5)
        left = Inches(0.25)
        width = Inches(9.25)
        height = Inches(5.0)
        # Flatten the pivot table by resetting the index
        # Create a table on the slide
        df_to_table(slide, report_data.xs(manager, level=0).reset_index(),
                    left, top, width, height)
    prs.save(output)


if __name__ == "__main__":

	con = None
	con = mdb.connect('localhost','root','','');
	cur = con.cursor()

	try:
		os.remove('/media/sf_transfer/' + outfile)
	except:
		print('No file\n')

	start = time.time()

	###############
	# Read in data
	###############
	query_eff = "SELECT parent_id,account_id,account_name,BEEs,target_audience,video_title,video_category,industry_name,min_time,video_id,effectiveness_absolute FROM %s.%s" % (DBNAME,TABLENAME_EFF)
	query_reach_act = "SELECT * FROM %s.%s WHERE cur_year = '1-Year'" % (DBNAME,TABLENAME_REACH_act)
	query_reach_act_2014 = "SELECT * FROM %s.%s WHERE cur_year = '1-Year'" % (DBNAME,TABLENAME_REACH_act_2014)
	query_reach_video = "SELECT * FROM %s.%s WHERE cur_year = '1-Year'" % (DBNAME,TABLENAME_REACH_video)
	query_reach_video_2014 = "SELECT * FROM %s.%s WHERE cur_year = '1-Year'" % (DBNAME,TABLENAME_REACH_video_2014)

	eff_df = pd.read_sql(query_eff,con)
	eff_df['effectiveness_absolute'] = eff_df['effectiveness_absolute'].astype(float)

	accounts = list(set(eff_df.account_id))

	for ppp in range(0,len(accounts)):

		printf('[create_ppt.py] Account %5s of %5s: G2 Account Id = %5s ... %.2f sec\n',ppp+1,len(accounts),accounts[ppp],time.time()-start)
		tmp_eff_df = eff_df[(eff_df['min_time'] > (cur_datetime.date()-timedelta(days=timeframe_days))) & (eff_df['account_id'] == str(accounts[ppp])) ]

		if (len(tmp_eff_df) > 0):
			###################################
			# Calculate Reach / Effectiveness 
			###################################
			eff_act_df = np.median(tmp_eff_df['effectiveness_absolute'])
			#eff_video_df = tmp_eff_df.groupby('video_id').agg([np.median])
			eff_video_df = tmp_eff_df.groupby('video_id').median().reset_index()

			cur_video_title = []
			for i in range(0,len(eff_video_df)):
				cur_video_title.append(tmp_eff_df['video_title'][tmp_eff_df.index[list(tmp_eff_df['video_id']).index(eff_video_df['video_id'][i])] ])
			eff_video_df = eff_video_df.join(pd.DataFrame(cur_video_title)).rename(columns={0:'video_title'}).sort('effectiveness_absolute',ascending=0)
			reach_act_df = pd.read_sql(query_reach_act,con)
			reach_act_2014_df = pd.read_sql(query_reach_act_2014,con)
			reach_video_df = pd.read_sql(query_reach_video,con)
			reach_video_2014_df = pd.read_sql(query_reach_video_2014,con)

			eff_video_df = eff_video_df[(eff_video_df['effectiveness_absolute'] == eff_video_df['effectiveness_absolute'])] ## Remove any NaN values 

			######################
			# Start PPT creation
			######################
			#args = parse_args()
			#infile = 'simple-template.pptx'	
			infile = 'qbr_template.pptx'	
			outfile = 'myreport.pptx'
 			report_name = './data/sales-funnel.xlsx' 
			chart = "report-image.png"

			prs = Presentation(infile)

			###################################################
			# Create Account/Video-Level Effectiveness Chart
			###################################################
			graph_slide_layout = prs.slide_layouts[0]
			slide = prs.slides.add_slide(graph_slide_layout)
			title = slide.shapes.title
			title.text = "Yearly Results (Video-Level)"

			#####################################
			# 1) Create MKTG Effectiveness Chart
			#####################################
			cur_reach_act_df = reach_act_df[(reach_act_df.account_id == int(accounts[ppp]) )]	
			cur_reach_act_df['USER_reach'] = [i if i < 1 else 1 for i in cur_reach_act_df['USER_reach']]	
			cur_reach_video_df = reach_video_df[(reach_video_df.account_id == int(accounts[ppp]) )].sort('USER_reach',ascending=0)	
			cur_reach_video_df['USER_reach'] = [i if i < 1 else 1 for i in cur_reach_video_df['USER_reach']]	
			cur_reach_act_2014_df = reach_act_2014_df[(reach_act_2014_df.account_id == int(accounts[ppp]) )]	
			cur_reach_act_2014_df['USER_reach'] = [i if i < 1 else 1 for i in cur_reach_act_2014_df['USER_reach']]	
			cur_reach_video_2014_df = reach_video_2014_df[(reach_video_2014_df.account_id == int(accounts[ppp]) )]	
			cur_reach_video_2014_df['USER_reach'] = [i if i < 1 else 1 for i in cur_reach_video_2014_df['USER_reach']]	
#	
#			###############
#			# Add Title 
#			###############
#			left,top,width,height = Inches(0.125),Inches(0.90),Inches(3.5),Inches(3)
#			#left = top = width = height = Inches(1)
#			shape = slide.shapes.add_shape(
#				#MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
#				MSO_SHAPE.RECTANGLE, left, top, width, height
#			)
#
#			fill = shape.fill
#			#fill.background() ## No Fill
#			fill.solid() ## No Fill
#			fill.fore_color.rgb = RGBColor(255,255,255)
#
#			line = shape.line
#			line.fill.background()
#
#			left,top,width,height = Inches(0.125),Inches(0.90),Inches(3.5),Inches(0.5)
#			#left = top = width = height = Inches(1)
#			shape2 = slide.shapes.add_shape(
#				MSO_SHAPE.RECTANGLE, left, top, width, height
#			)
#
#			fill = shape2.fill
#			fill.solid()
#			fill.fore_color.rgb = RGBColor(85,159,206)
#
#			line = shape2.line
#			line.fill.background()
#
#			text_frame = shape2.text_frame
#			text_frame.text = 'Marketing Effectiveness'
#			text_frame.margin_bottom = Inches(0.08)
#			text_frame.margin_left = 0
#			text_frame.vertical_anchor = MSO_ANCHOR.TOP
#			text_frame.word_wrap = False
#			#text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT	
#
#			# define chart data ---------------------
#			chart_data = ChartData()
#			chart_data.categories = cur_reach_video_df['video_title']
#			chart_data.add_series('Series 1', cur_reach_video_df['USER_reach']*100 )
#
#			#######################
#			# Add chart to slide 
#			#######################
#			x, y, cx, cy = Inches(0.125), Inches(1.35), Inches(3.5), Inches(2.5)
#			graphic_frame = slide.shapes.add_chart(
#				XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
#			)
#
#			####################
#			# Change Font
#			####################
#			chart = graphic_frame.chart
#
#			category_axis = chart.category_axis
#			#category_axis.has_major_gridlines = True
#			category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
#			category_axis.tick_labels.font.italic = True
#			category_axis.tick_labels.font.size = Pt(10)
#			category_axis.tick_labels.font.color.rgb = RGBColor(0,0,0)
#
#			value_axis = chart.value_axis
#			value_axis.maximum_scale = 100
#			value_axis.major_unit = 20
#			value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
#			value_axis.tick_labels.font.color.rgb = RGBColor(0,0,0) 
#			#value_axis.has_minor_gridlines = True
#
#			tick_labels = value_axis.tick_labels
#			tick_labels.number_format = '0"%"'
#			tick_labels.font.bold = True
#			tick_labels.font.size = Pt(14)
#
#			######################
#			# Add Data Labels
#			######################
#			if ADD_DATA_LABELS == True:
#				plot = chart.plots[0]
#				plot.has_data_labels = True
#				data_labels = plot.data_labels
#		
#				data_labels.font.size = Pt(13)
#				data_labels.number_format = '0.0"%"'
#				data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
#				data_labels.position = XL_LABEL_POSITION.INSIDE_END

			################################
			# 2) Create Video-Level Chart
			################################
			######################
			# Add Title 
			######################
			left,top,width,height = Inches(0.125),Inches(0.90),Inches(7.125),Inches(4.5)
			#left = top = width = height = Inches(1)
			shape = slide.shapes.add_shape(
				#MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
				MSO_SHAPE.RECTANGLE, left, top, width, height
			)
	
			fill = shape.fill
			#fill.background() ## No Fill
			fill.solid() ## No Fill
			fill.fore_color.rgb = RGBColor(255,255,255)

			line = shape.line
			line.fill.background()

			left,top,width,height = Inches(0.125),Inches(0.90),Inches(7.125),Inches(0.5)
			#left = top = width = height = Inches(1)
			shape2 = slide.shapes.add_shape(
				MSO_SHAPE.RECTANGLE, left, top, width, height
			)

			fill = shape2.fill
			fill.solid()
			fill.fore_color.rgb = RGBColor(85,159,206)

			line = shape2.line
			line.fill.background()

			text_frame = shape2.text_frame
			text_frame.text = 'Content Effectiveness'
			text_frame.margin_bottom = Inches(0.08)
			text_frame.margin_left = 0
			text_frame.vertical_anchor = MSO_ANCHOR.TOP
			text_frame.word_wrap = False
			#text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT	
			# define chart data ---------------------
			chart_data = ChartData()
			#chart_data.categories = ['East', 'West', 'Midwest']
			#test = {19.2, 21.4, 16.7}
			#chart_data.add_series('Series 1', test )

			eff_video_df['video_title'] = [x.decode('latin-1').encode('ascii','ignore').replace('&','') for x in eff_video_df['video_title']]
			chart_data.categories = eff_video_df['video_title']
			chart_data.add_series('Series 1', eff_video_df['effectiveness_absolute']*100 )

			#######################
			# Add chart to slide 
			#######################
			x, y, cx, cy = Inches(0.25), Inches(1.5), Inches(6.875), Inches(3.75)
			#x, y, cx, cy = Inches(0.125), Inches(1.35), Inches(3.5), Inches(2.5)
			try:
				graphic_frame = slide.shapes.add_chart(
					XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
				)

				####################
				# Change Font
				####################
				chart = graphic_frame.chart

				category_axis = chart.category_axis  # X-Axis
				#category_axis.has_major_gridlines = True
				category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
				category_axis.tick_labels.font.italic = True
				category_axis.tick_labels.font.size = Pt(10)
				category_axis.tick_labels.font.color.rgb = RGBColor(0,0,0)

				value_axis = chart.value_axis
				value_axis.maximum_scale = 100       # Y-Axis
				value_axis.major_unit = 20
				value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
				value_axis.tick_labels.font.color.rgb = RGBColor(0,0,0) 
				#value_axis.has_minor_gridlines = True

				tick_labels = value_axis.tick_labels
				tick_labels.number_format = '0"%"'
				tick_labels.font.bold = True
				tick_labels.font.size = Pt(14)

			except Exception as e:
				printf("ppp = %4d:Account_id = %s: Line %5s: %70s\n",ppp,accounts[ppp],sys.exc_traceback.tb_lineno,e)
				sys.exit() 

			######################
			# Add Data Labels
			######################
			if ADD_DATA_LABELS == True:
				plot = chart.plots[0]
				plot.has_data_labels = True
				data_labels = plot.data_labels

				data_labels.font.size = Pt(13)
				data_labels.number_format = '0.0"%"'
				data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
				data_labels.position = XL_LABEL_POSITION.INSIDE_END

			################################
			# 3) Create Comments Chart
			################################
			######################
			# Add Title 
			######################
			left,top,width,height = Inches(7.375),Inches(0.90),Inches(2.5),Inches(4.5)
			shape = slide.shapes.add_shape(
				MSO_SHAPE.RECTANGLE, left, top, width, height
			)

			fill = shape.fill
			#fill.background() ## No Fill
			fill.solid() ## Solid Fill 
			fill.fore_color.rgb = RGBColor(255,255,255)

			line = shape.line
			line.fill.background()

			text_frame = shape.text_frame
			text_frame.margin_bottom = Inches(0.08)
			text_frame.margin_left = Inches(0.20)
			text_frame.vertical_anchor = MSO_ANCHOR.TOP
			text_frame.word_wrap = False

			p = text_frame.paragraphs[0]
			#run = p.add_run()
			p.text = '\n\n\nComment 1\nComment 2\nComment 3'
			p.font.name = 'Calibri'
			p.font.size = Pt(14)
			p.alignment = PP_ALIGN.LEFT
			p.font.color.rgb = RGBColor(0,0,0) 

			############ Header #############
			left,top,width,height = Inches(7.375),Inches(0.90),Inches(2.5),Inches(0.5)
			#left = top = width = height = Inches(1)
			shape2 = slide.shapes.add_shape(
				MSO_SHAPE.RECTANGLE, left, top, width, height
			)

			fill = shape2.fill
			fill.solid()
			fill.fore_color.rgb = RGBColor(85,159,206)

			line = shape2.line
			line.fill.background()

			text_frame = shape2.text_frame
			text_frame.text = 'Comments'
			text_frame.margin_bottom = Inches(0.08)
			text_frame.margin_left = 0
			text_frame.vertical_anchor = MSO_ANCHOR.TOP
			text_frame.word_wrap = False

			###################################################
			# Create YoY Video-Level Slide 
			###################################################
			graph_slide_layout = prs.slide_layouts[0]
			slide = prs.slides.add_slide(graph_slide_layout)
			title = slide.shapes.title
			title.text = "Yearly Results (Video-Level)"
			#placeholder = slide.placeholders[1]
	
			########################
			# 1) Create YoY table
			########################
			rows,cols = len(cur_reach_video_df)+1,4
			top = Inches(0.875)
			left = Inches(0.375)
			width = Inches(9.25)
			height = Inches(0.6)
	
			table = slide.shapes.add_table(rows, cols, left, top, width, height).table
	
			# set column widths
			#table.columns[0].width = Inches(2.0)
			#table.columns[1].width = Inches(4.0)
	
			for j in range(0,cols):
				table.rows[0].cells[j].textframe.paragraphs[0].font.size = Pt(16)
				table.rows[0].cells[j].textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

			for j in range(0,cols):
				for i in range(1,rows):
					table.rows[i].cells[j].textframe.paragraphs[0].font.size = Pt(12)
					if (j > 0):
						table.rows[i].cells[j].textframe.paragraphs[0].alignment = PP_ALIGN.CENTER
					table.rows[i].cells[j].vertical_anchor = MSO_ANCHOR.MIDDLE
					table.rows[i].cells[j].margin_bottom = Inches(0.025) 
					table.rows[i].cells[j].margin_top = Inches(0.025) 

			# write column headings
			table.cell(0, 0).text = 'Video Title'
			table.cell(0, 1).text = 'Last 12-Months'
			table.cell(0, 2).text = 'Previous 12-Months'
			table.cell(0, 3).text = 'YoY Growth'
	
			# write body cells
			for i in range(0,len(cur_reach_video_df)):
				idx = cur_reach_video_df.index[i]
				table.cell(i+1, 0).text = cur_reach_video_df.ix[idx].video_title.decode('latin-1').encode('ascii','ignore') 
				table.cell(i+1, 1).text = str(cur_reach_video_df.ix[idx].Nparent) 
				try:
					idx_2014 = cur_reach_video_2014_df.index[all_indices(cur_reach_video_df.ix[idx].video_id,cur_reach_video_2014_df.video_id)[0]]
					table.cell(i+1, 2).text = str(cur_reach_video_2014_df.ix[idx_2014].Nparent) 
					table.cell(i+1, 3).text = str("%.1f%%" % (100. * (float(cur_reach_video_df.ix[idx].Nparent) / float(cur_reach_video_2014_df.ix[idx_2014].Nparent) - 1)) ) 
				except:
					table.cell(i+1, 2).text = 'NA' 
					table.cell(i+1, 3).text = 'NA' 

				#table.cell(i+1, 2).text = 'Qux'

			## Put total account values in
			#table.cell(len(cur_reach_video_df)+1, 0).text = 'Total' 
			#table.cell(len(cur_reach_video_df)+1, 1).text = str(cur_reach_act_df.Nparent) 
			#try:
			#	table.cell(i+1, 2).text = str(cur_reach_act_2014_df.Nparent) 
			#	table.cell(i+1, 3).text = str("%.1f%%" % (100. * (float(cur_reach_act_df.Nparent) / float(cur_reach_video_2014_df.Nparent) - 1)) ) 
			#except:
			#	table.cell(i+1, 2).text = 'NA' 
			#	table.cell(i+1, 3).text = 'NA' 

			######################
			# Add Title 
			######################
			left,top,width,height = Inches(0),Inches(4.625),Inches(10),Inches(1)
			shape = slide.shapes.add_shape(
				MSO_SHAPE.RECTANGLE, left, top, width, height
			)

			fill = shape.fill
			#fill.background() ## No Fill
			fill.solid() ## Solid Fill 
			fill.fore_color.rgb = RGBColor(255,255,255)

			line = shape.line
			line.fill.background()

			text_frame = shape.text_frame
			text_frame.margin_bottom = Inches(0.08)
			text_frame.margin_left = Inches(0.20)
			text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
			text_frame.word_wrap = False

			p = text_frame.paragraphs[0]
			#run = p.add_run()
			p.text = 'Comment 1\nComment 2'
			p.font.name = 'Calibri'
			p.font.size = Pt(14)
			p.alignment = PP_ALIGN.CENTER
			p.font.color.rgb = RGBColor(0,0,0) 

			###################################################
			# Create Survey Slide 
			###################################################
			graph_slide_layout = prs.slide_layouts[0]
			slide = prs.slides.add_slide(graph_slide_layout)
			title = slide.shapes.title
			title.text = "Employee Survey Review (8/1/15-present)"
			#placeholder = slide.placeholders[1]
			# Create a slide for each manager
	
			query_survey = "SELECT * FROM %s.%s WHERE account_id = %s" % (SANDBOX,TABLENAME_SURVEY,accounts[ppp])
			question1 = 'The video content was clear and easy to understand';
			question2 = 'The video addressed questions I had on this topic';
			question3 = 'I would recommend this video to';
			query_survey_question1 = "SELECT account_id,scale_value,COUNT(scale_value) as Nscale_cnt, \
									SUM(Nsurvey_response) as Nscale_survey_response \
									FROM %s.%s \
									WHERE account_id = %s AND question_text LIKE '%%%s%%' GROUP BY scale_value" \
									% (SANDBOX,TABLENAME_SURVEY,accounts[ppp],question1)

			query_survey_question2 = "SELECT account_id,scale_value,COUNT(scale_value) as Nscale_cnt, \
									SUM(Nsurvey_response) as Nscale_survey_response \
									FROM %s.%s \
									WHERE account_id = %s AND question_text LIKE '%%%s%%' GROUP BY scale_value" \
									% (SANDBOX,TABLENAME_SURVEY,accounts[ppp],question2)

			query_survey_question3 = "SELECT account_id,scale_value,COUNT(scale_value) as Nscale_cnt, \
									SUM(Nsurvey_response) as Nscale_survey_response \
									FROM %s.%s \
									WHERE account_id = %s AND question_text LIKE '%%%s%%' GROUP BY scale_value" \
									% (SANDBOX,TABLENAME_SURVEY,accounts[ppp],question3)

			survey_df = pd.read_sql(query_survey,con)
			question1_act_df = pd.read_sql(query_survey_question1,con)
			question2_act_df = pd.read_sql(query_survey_question2,con)
			question3_act_df = pd.read_sql(query_survey_question3,con)

			if (len(survey_df) > 0):
				######################
				# Add Title 
				######################
				left,top,width,height = Inches(0.125),Inches(0.90),Inches(4.75),Inches(4.5)
				shape = slide.shapes.add_shape(
					MSO_SHAPE.RECTANGLE, left, top, width, height
				)

				fill = shape.fill
				#fill.background() ## No Fill
				fill.solid() ## Solid Fill 
				fill.fore_color.rgb = RGBColor(255,255,255)
	
				line = shape.line
				line.fill.background()
	
				text_frame = shape.text_frame
				text_frame.margin_bottom = Inches(0.08)
				text_frame.margin_left = Inches(0.20)
				text_frame.vertical_anchor = MSO_ANCHOR.TOP
				text_frame.word_wrap = False
	
				p = text_frame.paragraphs[0]
				#run = p.add_run()
				p.font.name = 'Calibri'
				p.font.size = Pt(14)
				p.alignment = PP_ALIGN.LEFT
				p.font.color.rgb = RGBColor(0,0,0) 
	
				############ Header #############
				left,top,width,height = Inches(0.125),Inches(0.90),Inches(4.75),Inches(0.5)
				#left = top = width = height = Inches(1)
				shape2 = slide.shapes.add_shape(
					MSO_SHAPE.RECTANGLE, left, top, width, height
				)
	
				fill = shape2.fill
				fill.solid()
				fill.fore_color.rgb = RGBColor(85,159,206)
	
				line = shape2.line
				line.fill.background()
	
				text_frame = shape2.text_frame
				text_frame.text = 'Average Responses for all Titles in Library'
				text_frame.margin_bottom = Inches(0.08)
				text_frame.margin_left = 0
				text_frame.vertical_anchor = MSO_ANCHOR.TOP
				text_frame.word_wrap = False
	
				####################################
				# 1) Pie Chart #1 with Average Data
				####################################
				disagree = sum(question1_act_df[(question1_act_df.scale_value == 1) | (question1_act_df.scale_value == 2)].Nscale_survey_response)
				neutral = sum(question1_act_df[(question1_act_df.scale_value == 3)].Nscale_survey_response)
				agree = sum(question1_act_df[(question1_act_df.scale_value == 4) | (question1_act_df.scale_value == 5)].Nscale_survey_response)
				total1 = disagree + neutral + agree
				if (total1 > 0):
					chart_data = ChartData()
					chart_data.categories = ['Agree', 'Neutral', 'Disagree']
					chart_data.add_series('Series 1', (agree/total1, neutral/total1, disagree/total1))
	
					x, y, cx, cy = Inches(0.25), Inches(1.5), Inches(2.0), Inches(1.25)
					chart = slide.shapes.add_chart(
						XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
					).chart
		
					chart.has_legend = True
					chart.legend.position = XL_LEGEND_POSITION.RIGHT
					chart.legend.include_in_layout = False
					chart.legend.font.size = Pt(10)
					chart.legend.font.color.rgb = RGBColor(0, 0, 0)
	
					chart.plots[0].has_data_labels = True
					data_labels = chart.plots[0].data_labels
					data_labels.font.size = Pt(10)
					data_labels.number_format = '0%'
					#data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
					data_labels.font.color.rgb = RGBColor(0, 0, 0)
					data_labels.position = XL_LABEL_POSITION.INSIDE_END
	
				####################################
				# 2) Pie Chart #2 with Average Data
				####################################
				disagree = sum(question2_act_df[(question2_act_df.scale_value == 1) | (question2_act_df.scale_value == 2)].Nscale_survey_response)
				neutral = sum(question2_act_df[(question2_act_df.scale_value == 3)].Nscale_survey_response)
				agree = sum(question2_act_df[(question2_act_df.scale_value == 4) | (question2_act_df.scale_value == 5)].Nscale_survey_response)
				total2 = disagree + neutral + agree
				if (total2 > 0):
					chart_data = ChartData()
					chart_data.categories = ['Agree', 'Neutral', 'Disagree']
					chart_data.add_series('Series 1', (agree/total2, neutral/total2, disagree/total2))
	
					x, y, cx, cy = Inches(0.25), Inches(2.75), Inches(2.0), Inches(1.25)
					chart = slide.shapes.add_chart(
						XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
					).chart

					chart.has_legend = True
					chart.legend.position = XL_LEGEND_POSITION.RIGHT
					chart.legend.include_in_layout = False
					chart.legend.font.size = Pt(10)
					chart.legend.font.color.rgb = RGBColor(0, 0, 0)

					chart.plots[0].has_data_labels = True
					data_labels = chart.plots[0].data_labels
					data_labels.font.size = Pt(10)
					data_labels.number_format = '0%'
					#data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
					data_labels.font.color.rgb = RGBColor(0, 0, 0)
					data_labels.position = XL_LABEL_POSITION.INSIDE_END
	
				####################################
				# 3) Pie Chart #3 with Average Data
				####################################
				disagree = sum(question3_act_df[(question3_act_df.scale_value == 1) | (question3_act_df.scale_value == 2)].Nscale_survey_response)
				neutral = sum(question3_act_df[(question3_act_df.scale_value == 3)].Nscale_survey_response)
				agree = sum(question3_act_df[(question3_act_df.scale_value == 4) | (question3_act_df.scale_value == 5)].Nscale_survey_response)
				total3 = disagree + neutral + agree
				if (total3 > 0):
					chart_data = ChartData()
					chart_data.categories = ['Agree', 'Neutral', 'Disagree']
					chart_data.add_series('Series 1', (agree/total3, neutral/total3, disagree/total3))

					x, y, cx, cy = Inches(0.25), Inches(4.0), Inches(2.0), Inches(1.25)
					chart = slide.shapes.add_chart(
						XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
					).chart

					chart.has_legend = True
					chart.legend.position = XL_LEGEND_POSITION.RIGHT
					chart.legend.include_in_layout = False
					chart.legend.font.size = Pt(10)
					chart.legend.font.color.rgb = RGBColor(0, 0, 0)

					chart.plots[0].has_data_labels = True
					data_labels = chart.plots[0].data_labels
					data_labels.font.size = Pt(10)
					data_labels.number_format = '0%'
					#data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
					data_labels.font.color.rgb = RGBColor(0, 0, 0)
					data_labels.position = XL_LABEL_POSITION.INSIDE_END

				######################
				# Add Input Text 
				######################
				text_frame = shape.text_frame
				text_frame.margin_bottom = Inches(0.08)
				text_frame.margin_left = Inches(2.5)
				text_frame.vertical_anchor = MSO_ANCHOR.TOP
				text_frame.word_wrap = False
	
				question1_out = 'The video content was clear\nand easy to understand\n(' + str(total1) + ' responses)';
				question2_out = 'The video addressed questions\nI had on this topic\n(' + str(total2) + ' responses)';
				question3_out = 'I would recommend this video\nto colleagues\n(' + str(total3) + ' responses)';

				p = text_frame.paragraphs[0]
				#run = p.add_run()
				p.text = '\n\n\n\n' + question1_out + '\n\n\n\n' + question2_out + '\n\n\n\n' + question3_out  
				p.font.name = 'Calibri'
				p.font.size = Pt(12)
				p.alignment = PP_ALIGN.LEFT
				p.font.color.rgb = RGBColor(0,0,0) 

				######################
				# Add Input Text 
				######################
				left,top,width,height = Inches(5),Inches(0.9),Inches(4.875),Inches(4.5)
				shape = slide.shapes.add_shape(
					MSO_SHAPE.RECTANGLE, left, top, width, height
				)

				fill = shape.fill
				#fill.background() ## No Fill
				fill.solid() ## Solid Fill 
				fill.fore_color.rgb = RGBColor(255,255,255)
	
				line = shape.line
				line.fill.background()

				text_frame = shape.text_frame
				text_frame.margin_bottom = Inches(0.08)
				text_frame.margin_left = Inches(0.20)
				text_frame.vertical_anchor = MSO_ANCHOR.TOP
				text_frame.word_wrap = False

				p = text_frame.paragraphs[0]
				#run = p.add_run()
				p.text = '\n\n\nComment 1\nComment 2\nComment 3'
				p.font.name = 'Calibri'
				p.font.size = Pt(14)
				p.alignment = PP_ALIGN.LEFT
				p.font.color.rgb = RGBColor(0,0,0) 

				############ Header #############
				left,top,width,height = Inches(5),Inches(0.9),Inches(4.875),Inches(0.5)
				#left = top = width = height = Inches(1)
				shape2 = slide.shapes.add_shape(
					MSO_SHAPE.RECTANGLE, left, top, width, height
				)

				fill = shape2.fill
				fill.solid()
				fill.fore_color.rgb = RGBColor(85,159,206)

				line = shape2.line
				line.fill.background()

				text_frame = shape2.text_frame
				text_frame.text = 'Comments'
				text_frame.margin_bottom = Inches(0.08)
				text_frame.margin_left = 0
				text_frame.vertical_anchor = MSO_ANCHOR.TOP
				text_frame.word_wrap = False

			try:	
				cur_account_name = list(set(tmp_eff_df.account_name))[0].replace('/','').replace('\x92','')
				file_out = './csm_reports/act' + str(accounts[ppp]).zfill(5) + '_' + cur_account_name + '_results.pptx'
			except:
				file_out = './csm_reports/act' + str(accounts[ppp]).zfill(5) + '_results.pptx'
			
			prs.save(file_out.replace(' ',''))


